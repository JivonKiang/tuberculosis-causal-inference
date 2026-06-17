#!/usr/bin/env python3
"""
Generate A4 PPTX with embedded figures and tables.
Tables also exported as XLSX and DOCX.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import pandas as pd
import numpy as np
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, Border, Side, PatternFill, numbers
from openpyxl.utils import get_column_letter
from docx import Document
from docx.shared import Pt as DocxPt, Inches as DocxInches, Cm as DocxCm, RGBColor as DocxRGB
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH

OUTPUT_DIR = Path(r'E:\20260617 tuberculosis-causal-inference\04_图表')
DATA_PATH = Path(r'E:\20260617 tuberculosis-causal-inference\02_原始数据\who_gbd\master_data.csv')

REGION_NAMES = {
    'AFR': 'African', 'AMR': 'Americas', 'EMR': 'Eastern Mediterranean',
    'EUR': 'European', 'SEAR': 'South-East Asia', 'WPR': 'Western Pacific',
}

# A4 portrait dimensions in inches
A4_W = 8.27  # 210mm
A4_H = 11.69  # 297mm

# Load data
df = pd.read_csv(DATA_PATH)
df_valid = df[df['gbd_inc'] > 1.5].copy()

# ============================================================
# PART 1: Generate PPTX
# ============================================================
print("Creating A4 PPTX...")
prs = Presentation()
prs.slide_width = Cm(21.0)
prs.slide_height = Cm(29.7)

# Blank slide layout
blank_layout = prs.slide_layouts[6]  # blank

def add_title(slide, text, top=Cm(0.3), left=Cm(1.0), width=Cm(19.0), fontsize=18):
    txBox = slide.shapes.add_textbox(left, top, width, Cm(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(fontsize)
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(0x1a, 0x3a, 0x6b)

def add_image(slide, img_path, left, top, width=None, height=None):
    if width:
        slide.shapes.add_picture(str(img_path), left, top, width=width)
    elif height:
        slide.shapes.add_picture(str(img_path), left, top, height=height)

def add_table(slide, df_table, left, top, width, col_widths=None, fontsize=10):
    rows, cols = df_table.shape
    table_shape = slide.shapes.add_table(rows + 1, cols, left, top, width, Cm(0.6 * (rows + 1)))
    table = table_shape.table

    # Header
    for j, col_name in enumerate(df_table.columns):
        cell = table.cell(0, j)
        cell.text = str(col_name)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(fontsize)
            paragraph.font.name = 'Arial'
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1a, 0x3a, 0x6b)

    # Data
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i + 1, j)
            val = df_table.iloc[i, j]
            if isinstance(val, float):
                if abs(val) < 0.01 or abs(val) > 10000:
                    cell.text = f'{val:.2e}'
                else:
                    cell.text = f'{val:.1f}' if val == int(val) or abs(val) >= 1 else f'{val:.2f}'
            else:
                cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(fontsize - 1)
                paragraph.font.name = 'Arial'
                paragraph.alignment = PP_ALIGN.CENTER
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)

    return table_shape


# --- Slide 1: Title Slide ---
slide = prs.slides.add_slide(blank_layout)
txBox = slide.shapes.add_textbox(Cm(2), Cm(8), Cm(17), Cm(5))
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = 'Global Divergence in Tuberculosis\nIncidence Estimates'
p.font.size = Pt(32)
p.font.name = 'Arial'
p.font.color.rgb = RGBColor(0x1a, 0x3a, 0x6b)
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = 'WHO Global TB Report vs. GBD 2021: A Comparative Analysis'
p2.font.size = Pt(18)
p2.font.name = 'Arial'
p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(16)

p3 = tf.add_paragraph()
p3.text = 'Figure Set for Military Medical Research (IF 22.9)'
p3.font.size = Pt(14)
p3.font.name = 'Arial'
p3.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(40)

p4 = tf.add_paragraph()
p4.text = f'217 countries · {len(df_valid)} with valid GBD estimates · Data: WHO 2025, GBD 2021'
p4.font.size = Pt(12)
p4.font.name = 'Arial'
p4.font.color.rgb = RGBColor(0xaa, 0xaa, 0xaa)
p4.alignment = PP_ALIGN.CENTER
p4.space_before = Pt(20)

# --- Slide 2: Figure 1 - Global Map ---
slide = prs.slides.add_slide(blank_layout)
add_title(slide, 'Figure 1 | Global Divergence in TB Incidence Estimates')
img_path = OUTPUT_DIR / 'Fig1_Global_Divergence.png'
if img_path.exists():
    add_image(slide, img_path, left=Cm(1.5), top=Cm(2.0), width=Cm(18.0))

# --- Slide 3: Figure 2 - Bland-Altman ---
slide = prs.slides.add_slide(blank_layout)
add_title(slide, 'Figure 2 | Bland-Altman Decomposition & Regional Heterogeneity')
img_path = OUTPUT_DIR / 'Fig2_Bland_Altman.png'
if img_path.exists():
    add_image(slide, img_path, left=Cm(0.8), top=Cm(2.0), width=Cm(19.5))

# --- Slide 4: Figure 3 - Regional Forest ---
slide = prs.slides.add_slide(blank_layout)
add_title(slide, 'Figure 3 | Regional Heterogeneity in TB Incidence Divergence')
img_path = OUTPUT_DIR / 'Fig3_Regional_Forest.png'
if img_path.exists():
    add_image(slide, img_path, left=Cm(1.0), top=Cm(2.0), width=Cm(19.0))

# --- Slide 5: Figure 4 - Top Countries ---
slide = prs.slides.add_slide(blank_layout)
add_title(slide, 'Figure 4 | Countries with Largest WHO-GBD Divergence')
img_path = OUTPUT_DIR / 'Fig4_Top_Countries.png'
if img_path.exists():
    add_image(slide, img_path, left=Cm(0.8), top=Cm(2.0), width=Cm(19.5))

# --- Slide 6: Figure 5 - Mean vs Ratio ---
slide = prs.slides.add_slide(blank_layout)
add_title(slide, 'Figure 5 | Burden-Dependent Divergence Pattern')
img_path = OUTPUT_DIR / 'Fig5_Mean_vs_Ratio.png'
if img_path.exists():
    add_image(slide, img_path, left=Cm(1.5), top=Cm(2.0), width=Cm(18.0))

# --- Slide 7: Figure 6 - Dashboard ---
slide = prs.slides.add_slide(blank_layout)
add_title(slide, 'Figure 6 | Summary Dashboard')
img_path = OUTPUT_DIR / 'Fig6_Dashboard.png'
if img_path.exists():
    add_image(slide, img_path, left=Cm(0.5), top=Cm(1.8), width=Cm(20.0))

# --- Slide 8: Table 1 - Regional Summary ---
slide = prs.slides.add_slide(blank_layout)
add_title(slide, 'Table 1 | Regional Summary Statistics of TB Incidence Divergence')

region_summary = df_valid.groupby('region').agg(
    N=('country', 'count'),
    WHO_Mean=('who_inc', 'mean'),
    WHO_SD=('who_inc', 'std'),
    GBD_Mean=('gbd_inc', 'mean'),
    GBD_SD=('gbd_inc', 'std'),
    Ratio_Mean=('ratio', 'mean'),
    Ratio_SD=('ratio', 'std'),
    Diff_Mean=('diff', 'mean'),
    Diff_Median=('diff', 'median'),
    WHO_gt_GBD=('ratio', lambda x: (x > 1).sum()),
    GBD_gt_WHO=('ratio', lambda x: (x < 1).sum()),
).reset_index()
region_summary['region'] = region_summary['region'].map(REGION_NAMES)

# Round
for col in ['WHO_Mean', 'WHO_SD', 'GBD_Mean', 'GBD_SD', 'Diff_Mean', 'Diff_Median']:
    region_summary[col] = region_summary[col].round(1)
region_summary['Ratio_Mean'] = region_summary['Ratio_Mean'].round(2)
region_summary['Ratio_SD'] = region_summary['Ratio_SD'].round(2)

# Rename for display
region_summary.columns = [
    'Region', 'N', 'WHO Mean', 'WHO SD', 'GBD Mean', 'GBD SD',
    'Ratio Mean', 'Ratio SD', 'Diff Mean', 'Diff Median', 'WHO>GBD', 'GBD>WHO'
]

add_table(slide, region_summary, left=Cm(0.5), top=Cm(2.0), width=Cm(20.0), fontsize=9)

# --- Slide 9: Table 2 - Top 30 Countries ---
slide = prs.slides.add_slide(blank_layout)
add_title(slide, 'Table 2 | Top 30 Countries by WHO-GBD Incidence Divergence')

top30 = df_valid.nlargest(30, 'diff')[['country', 'region', 'who_inc', 'gbd_inc', 'ratio', 'diff', 'mean_inc']].copy()
top30['region'] = top30['region'].map(REGION_NAMES)
top30.columns = ['Country', 'Region', 'WHO Inc', 'GBD Inc', 'Ratio', 'Difference', 'Mean Inc']
top30['WHO Inc'] = top30['WHO Inc'].round(1)
top30['GBD Inc'] = top30['GBD Inc'].round(1)
top30['Ratio'] = top30['Ratio'].round(2)
top30['Difference'] = top30['Difference'].round(1)
top30['Mean Inc'] = top30['Mean Inc'].round(1)

add_table(slide, top30.reset_index(drop=True), left=Cm(0.3), top=Cm(2.0), width=Cm(20.4), fontsize=8)

# --- Slide 10: Table 3 - Bottom 30 Countries ---
slide = prs.slides.add_slide(blank_layout)
add_title(slide, 'Table 3 | Bottom 30 Countries (GBD > WHO) by Incidence Divergence')

bottom30 = df_valid.nsmallest(30, 'diff')[['country', 'region', 'who_inc', 'gbd_inc', 'ratio', 'diff', 'mean_inc']].copy()
bottom30['region'] = bottom30['region'].map(REGION_NAMES)
bottom30.columns = ['Country', 'Region', 'WHO Inc', 'GBD Inc', 'Ratio', 'Difference', 'Mean Inc']
bottom30['WHO Inc'] = bottom30['WHO Inc'].round(1)
bottom30['GBD Inc'] = bottom30['GBD Inc'].round(1)
bottom30['Ratio'] = bottom30['Ratio'].round(2)
bottom30['Difference'] = bottom30['Difference'].round(1)
bottom30['Mean Inc'] = bottom30['Mean Inc'].round(1)

add_table(slide, bottom30.reset_index(drop=True), left=Cm(0.3), top=Cm(2.0), width=Cm(20.4), fontsize=8)

pptx_path = OUTPUT_DIR / 'TB_Incidence_Divergence_Figures_Tables.pptx'
prs.save(str(pptx_path))
print(f"  PPTX saved: {pptx_path}")

# ============================================================
# PART 2: Generate XLSX Tables
# ============================================================
print("\nCreating XLSX tables...")

wb = Workbook()
thin_border = Border(
    left=Side(style='thin', color='CCCCCC'),
    right=Side(style='thin', color='CCCCCC'),
    top=Side(style='thin', color='CCCCCC'),
    bottom=Side(style='thin', color='CCCCCC'),
)
header_fill = PatternFill(start_color='1A3A6B', end_color='1A3A6B', fill_type='solid')
header_font = Font(name='Arial', size=12, bold=True, color='FFFFFF')
data_font = Font(name='Arial', size=11)
alt_fill = PatternFill(start_color='F0F4F8', end_color='F0F4F8', fill_type='solid')

# Sheet 1: Regional Summary
ws1 = wb.active
ws1.title = 'Regional Summary'
cols_r = ['Region', 'N', 'WHO Mean', 'WHO SD', 'GBD Mean', 'GBD SD',
          'Ratio Mean', 'Ratio SD', 'Diff Mean', 'Diff Median', 'WHO>GBD', 'GBD>WHO']
for j, col_name in enumerate(cols_r, 1):
    cell = ws1.cell(row=1, column=j, value=col_name)
    cell.fill = header_fill
    cell.font = header_font
    cell.alignment = Alignment(horizontal='center', vertical='center')
    cell.border = thin_border
    ws1.column_dimensions[get_column_letter(j)].width = 14

for i, (_, row) in enumerate(region_summary.iterrows()):
    for j, val in enumerate(row):
        cell = ws1.cell(row=i + 2, column=j + 1, value=val)
        cell.font = data_font
        cell.border = thin_border
        cell.alignment = Alignment(horizontal='center')
        if i % 2 == 0:
            cell.fill = alt_fill

# Sheet 2: Top WHO>GBD
ws2 = wb.create_sheet('WHO > GBD (Top 30)')
cols_t = ['Country', 'Region', 'WHO Incidence', 'GBD Incidence', 'Ratio', 'Difference', 'Mean Incidence']
for j, col_name in enumerate(cols_t, 1):
    cell = ws2.cell(row=1, column=j, value=col_name)
    cell.fill = header_fill
    cell.font = header_font
    cell.alignment = Alignment(horizontal='center', vertical='center')
    cell.border = thin_border
    ws2.column_dimensions[get_column_letter(j)].width = 16

top30_full = df_valid.nlargest(30, 'diff')[['country', 'region', 'who_inc', 'gbd_inc', 'ratio', 'diff', 'mean_inc']]
for i, (_, row) in enumerate(top30_full.iterrows()):
    vals = [row['country'], REGION_NAMES.get(row['region'], row['region']),
            round(row['who_inc'], 1), round(row['gbd_inc'], 1),
            round(row['ratio'], 2), round(row['diff'], 1), round(row['mean_inc'], 1)]
    for j, val in enumerate(vals):
        cell = ws2.cell(row=i + 2, column=j + 1, value=val)
        cell.font = data_font
        cell.border = thin_border
        cell.alignment = Alignment(horizontal='center')
        if i % 2 == 0:
            cell.fill = alt_fill

# Sheet 3: Bottom GBD>WHO
ws3 = wb.create_sheet('GBD > WHO (Bottom 30)')
for j, col_name in enumerate(cols_t, 1):
    cell = ws3.cell(row=1, column=j, value=col_name)
    cell.fill = header_fill
    cell.font = header_font
    cell.alignment = Alignment(horizontal='center', vertical='center')
    cell.border = thin_border
    ws3.column_dimensions[get_column_letter(j)].width = 16

bottom30_full = df_valid.nsmallest(30, 'diff')[['country', 'region', 'who_inc', 'gbd_inc', 'ratio', 'diff', 'mean_inc']]
for i, (_, row) in enumerate(bottom30_full.iterrows()):
    vals = [row['country'], REGION_NAMES.get(row['region'], row['region']),
            round(row['who_inc'], 1), round(row['gbd_inc'], 1),
            round(row['ratio'], 2), round(row['diff'], 1), round(row['mean_inc'], 1)]
    for j, val in enumerate(vals):
        cell = ws3.cell(row=i + 2, column=j + 1, value=val)
        cell.font = data_font
        cell.border = thin_border
        cell.alignment = Alignment(horizontal='center')
        if i % 2 == 0:
            cell.fill = alt_fill

# Sheet 4: All Countries
ws4 = wb.create_sheet('All Countries (n=211)')
cols_all = ['Country', 'Region', 'WHO Incidence', 'GBD Incidence', 'Ratio', 'Difference', 'Mean', 'Lat', 'Lon']
for j, col_name in enumerate(cols_all, 1):
    cell = ws4.cell(row=1, column=j, value=col_name)
    cell.fill = header_fill
    cell.font = header_font
    cell.alignment = Alignment(horizontal='center', vertical='center')
    cell.border = thin_border
    ws4.column_dimensions[get_column_letter(j)].width = 15

for i, (_, row) in enumerate(df_valid.sort_values('diff', ascending=False, key=abs).iterrows()):
    vals = [row['country'], REGION_NAMES.get(row['region'], row['region']),
            round(row['who_inc'], 1), round(row['gbd_inc'], 1),
            round(row['ratio'], 2), round(row['diff'], 1), round(row['mean_inc'], 1),
            round(row['lat'], 1), round(row['lon'], 1)]
    for j, val in enumerate(vals):
        cell = ws4.cell(row=i + 2, column=j + 1, value=val)
        cell.font = data_font
        cell.border = thin_border
        cell.alignment = Alignment(horizontal='center')
        if i % 2 == 0:
            cell.fill = alt_fill

xlsx_path = OUTPUT_DIR / 'Table_TB_Divergence_Statistics.xlsx'
wb.save(str(xlsx_path))
print(f"  XLSX saved: {xlsx_path}")

# ============================================================
# PART 3: Generate DOCX Tables
# ============================================================
print("\nCreating DOCX tables...")

doc = Document()

# Set default font
style = doc.styles['Normal']
style.font.name = 'Arial'
style.font.size = DocxPt(11)
style.paragraph_format.space_after = DocxPt(4)
style.paragraph_format.space_before = DocxPt(2)

# Title
title = doc.add_heading('TB Incidence Divergence: WHO vs GBD', level=1)
title_run = title.runs[0]
title_run.font.name = 'Arial'
title_run.font.size = DocxPt(16)
title_run.font.color.rgb = DocxRGB(0x1a, 0x3a, 0x6b)

doc.add_paragraph(f'Data Source: WHO Global TB Report 2025; GBD 2021 estimates.', style='Normal')
doc.add_paragraph(f'Countries analyzed: {len(df_valid)} (filtered from 217 records).', style='Normal')
doc.add_paragraph()

# Table 1: Regional Summary
doc.add_heading('Table 1. Regional Summary Statistics', level=2)
doc.add_paragraph(
    'Summary of WHO and GBD TB incidence estimates by WHO region. '
    'Values are mean ± SD of country-level incidence rates per 100,000 population.',
    style='Normal'
)

table1 = doc.add_table(rows=len(region_summary) + 1, cols=len(region_summary.columns))
table1.style = 'Light Grid Accent 1'
table1.alignment = WD_TABLE_ALIGNMENT.CENTER

for j, col_name in enumerate(region_summary.columns):
    cell = table1.cell(0, j)
    cell.text = col_name
    for paragraph in cell.paragraphs:
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in paragraph.runs:
            run.font.size = DocxPt(9)
            run.font.bold = True
            run.font.name = 'Arial'

for i, (_, row) in enumerate(region_summary.iterrows()):
    for j, val in enumerate(row):
        cell = table1.cell(i + 1, j)
        cell.text = str(val)
        for paragraph in cell.paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in paragraph.runs:
                run.font.size = DocxPt(9)
                run.font.name = 'Arial'

doc.add_paragraph()

# Table 2: Top 30 WHO > GBD
doc.add_heading('Table 2. Top 30 Countries (WHO > GBD)', level=2)
doc.add_paragraph('Countries where WHO estimates exceed GBD estimates by the largest absolute margin.', style='Normal')

top30_cols = ['Country', 'Region', 'WHO Inc', 'GBD Inc', 'Ratio', 'Difference', 'Mean Inc']
table2 = doc.add_table(rows=31, cols=7)
table2.style = 'Light Grid Accent 1'
table2.alignment = WD_TABLE_ALIGNMENT.CENTER

for j, col_name in enumerate(top30_cols):
    cell = table2.cell(0, j)
    cell.text = col_name
    for paragraph in cell.paragraphs:
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in paragraph.runs:
            run.font.size = DocxPt(9)
            run.font.bold = True
            run.font.name = 'Arial'

top30_data = df_valid.nlargest(30, 'diff')
for i, (_, row) in enumerate(top30_data.iterrows()):
    vals = [row['country'], REGION_NAMES.get(row['region'], row['region']),
            f"{row['who_inc']:.1f}", f"{row['gbd_inc']:.1f}",
            f"{row['ratio']:.2f}", f"{row['diff']:.1f}", f"{row['mean_inc']:.1f}"]
    for j, val in enumerate(vals):
        cell = table2.cell(i + 1, j)
        cell.text = val
        for paragraph in cell.paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in paragraph.runs:
                run.font.size = DocxPt(9)
                run.font.name = 'Arial'

doc.add_paragraph()

# Table 3: Bottom 30 GBD > WHO
doc.add_heading('Table 3. Bottom 30 Countries (GBD > WHO)', level=2)
doc.add_paragraph('Countries where GBD estimates exceed WHO estimates by the largest absolute margin.', style='Normal')

table3 = doc.add_table(rows=31, cols=7)
table3.style = 'Light Grid Accent 1'
table3.alignment = WD_TABLE_ALIGNMENT.CENTER

for j, col_name in enumerate(top30_cols):
    cell = table3.cell(0, j)
    cell.text = col_name
    for paragraph in cell.paragraphs:
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in paragraph.runs:
            run.font.size = DocxPt(9)
            run.font.bold = True
            run.font.name = 'Arial'

bottom30_data = df_valid.nsmallest(30, 'diff')
for i, (_, row) in enumerate(bottom30_data.iterrows()):
    vals = [row['country'], REGION_NAMES.get(row['region'], row['region']),
            f"{row['who_inc']:.1f}", f"{row['gbd_inc']:.1f}",
            f"{row['ratio']:.2f}", f"{row['diff']:.1f}", f"{row['mean_inc']:.1f}"]
    for j, val in enumerate(vals):
        cell = table3.cell(i + 1, j)
        cell.text = val
        for paragraph in cell.paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in paragraph.runs:
                run.font.size = DocxPt(9)
                run.font.name = 'Arial'

docx_path = OUTPUT_DIR / 'Table_TB_Divergence_Statistics.docx'
doc.save(str(docx_path))
print(f"  DOCX saved: {docx_path}")

print("\nAll files generated successfully!")
print(f"  PPTX: {pptx_path}")
print(f"  XLSX: {xlsx_path}")
print(f"  DOCX: {docx_path}")
